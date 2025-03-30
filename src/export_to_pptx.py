from pptx import Presentation
from pptx.util import Inches
from pathlib import Path

pptx_dir = Path("pptx_exports")
pptx_dir.mkdir(exist_ok=True)

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank slide
img_path = Path("logs_v2/plots/combined_figure.png")

# Insert image
left = Inches(0.5)
top = Inches(0.5)
height = Inches(6.5)
slide.shapes.add_picture(str(img_path), left, top, height=height)

prs.save(pptx_dir / "peptide_membrane_summary.pptx")
print("[✔] PowerPoint slide exported to pptx_exports/")

